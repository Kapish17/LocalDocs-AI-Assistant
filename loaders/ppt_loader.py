from langchain_core.documents import Document
from pptx import Presentation


def load_ppt(file_path):
    """
    Reads a PowerPoint file and extracts text.
    """

    presentation = Presentation(file_path)

    text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return Document(
        page_content=text,
        metadata={
            "source": file_path.name,
            "file_type": "pptx"
        }
    )